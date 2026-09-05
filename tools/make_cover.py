#!/usr/bin/env python3
"""Build the GridScript cover page (PPTX) and export to PDF via LibreOffice."""
import os
import subprocess

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
PPTX = os.path.join(ROOT, 'report', 'cover.pptx')
PDFDIR = os.path.join(ROOT, 'report')

FONT = 'DejaVu Sans'
BG = RGBColor(0x11, 0x18, 0x27)
ACCENT = RGBColor(0x22, 0xD3, 0xEE)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
ROW_ALT = RGBColor(0x1E, 0x29, 0x3B)
HEADER_BG = RGBColor(0x0E, 0x74, 0x90)

MEMBERS = [
    ('Member', 'Reg. number', 'Contribution'),
    ('Nabil Ismail Abdulkadir', 'UG22CSC1047', 'Interpreter & semantics: src/interpreter.py, src/builtins.py'),
    ('Ahmad Auwal Abubakar', 'UG22CSC1075', 'Parser & grammar: src/parser.py, BNF rules'),
    ('Abubakar Muhammad Sulaiman', 'UG22CSC1046', 'Lexer & tokens: src/lexer.py, regex rules'),
    ('Rukayya Musbahu Imam', 'UG22CSC1040', 'AST & scoping model: scope chains, static scoping'),
    ('Muhammad Salisu', 'UG20CSC1005', 'Tests & report: 61 unit + 17 integration tests'),
]

STAGES = [('LEXER', 'tokens'), ('PARSER', 'AST'), ('EVALUATOR', 'runs rules')]


def textbox(slide, left, top, width, height, text, size, color, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.name = FONT
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    textbox(slide, 0.7, 0.25, 11.9, 0.9, 'GridScript', 44, WHITE, bold=True)
    textbox(slide, 0.7, 1.15, 11.9, 0.5, 'A small interpreter for a 2D game actor',
            20, MUTED)

    # pipeline: three boxes with arrows
    x, box_w, gap = 0.7, 3.3, 1.3
    for i, (stage, sub) in enumerate(STAGES):
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(x), Inches(1.95), Inches(box_w), Inches(1.0))
        shp.fill.solid()
        shp.fill.fore_color.rgb = ROW_ALT
        shp.line.color.rgb = ACCENT
        shp.line.width = Pt(2)
        tf = shp.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = stage
        p.font.size = Pt(18)
        p.font.name = FONT
        p.font.bold = True
        p.font.color.rgb = ACCENT
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = sub
        p2.font.size = Pt(12)
        p2.font.name = FONT
        p2.font.color.rgb = MUTED
        p2.alignment = PP_ALIGN.CENTER
        if i < 2:
            textbox(slide, x + box_w + 0.25, 2.15, 0.8, 0.6, '\u2192', 36, ACCENT,
                    align=PP_ALIGN.CENTER)
        x += box_w + gap

    # members table
    shape = slide.shapes.add_table(6, 3, Inches(0.7), Inches(3.25), Inches(11.9), Inches(3.1))
    table = shape.table
    table.columns[0].width = Inches(3.6)
    table.columns[1].width = Inches(2.1)
    table.columns[2].width = Inches(6.2)
    for r, row in enumerate(MEMBERS):
        table.rows[r].height = Inches(0.5)
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = ''
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = HEADER_BG if r == 0 else (ROW_ALT if r % 2 else BG)
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.font.size = Pt(12 if r == 0 else 11)
            p.font.name = FONT
            p.font.bold = (r == 0)
            p.font.color.rgb = WHITE

    textbox(slide, 0.7, 6.6, 11.9, 0.4,
            'CSC4207 Organization of Programming Languages  |  Group 1  |  2026',
            11, MUTED, align=PP_ALIGN.CENTER)

    prs.save(PPTX)
    print('wrote', PPTX)


def to_pdf():
    subprocess.run(['soffice', '--headless', '--convert-to', 'pdf',
                    '--outdir', PDFDIR, PPTX],
                   capture_output=True, timeout=180, check=False)
    print('converted to', os.path.join(PDFDIR, 'cover.pdf'))


if __name__ == '__main__':
    build()
    to_pdf()
